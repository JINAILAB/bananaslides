from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation

from bananaslides.config import PipelineConfig
from bananaslides.domain.models import DetectionBox, ElementCategory, OCRLine, OCRResult, SlideSize
from bananaslides.modules.inpainting.base import Inpainter
from bananaslides.modules.ocr.base import OcrEngine
from bananaslides.modules.ppt.render import PowerPointRenderer
from bananaslides.modules.typesetting.fixed_font_typesetter import FixedFontTypesetter
from bananaslides.pipeline.orchestrator import PipelineComponents, SlideToPptPipeline


class _FakeDetector:
    backend_name = "fake_detector"

    def detect(self, slide_image_path: Path, slide_size: SlideSize) -> list[DetectionBox]:
        assert slide_image_path.exists()
        assert slide_size.width_px == 1200
        return [
            DetectionBox(
                box_id="t0001",
                category=ElementCategory.TEXT,
                x=100,
                y=100,
                width=400,
                height=80,
                confidence=0.95,
                label="plain text",
            )
        ]


class _FakeOcrEngine(OcrEngine):
    backend_name = "fake_ocr"

    def recognize(self, slide_image_path: Path, detections: list[DetectionBox]) -> list[OCRResult]:
        assert slide_image_path.exists()
        assert detections[0].box_id == "t0001"
        return [
            OCRResult(
                box_id="t0001",
                text="Hello PPT",
                lines=[
                    OCRLine(
                        text="Hello PPT",
                        bbox=[[10.0, 10.0], [90.0, 10.0], [90.0, 34.0], [10.0, 34.0]],
                        confidence=0.99,
                    )
                ],
                language="ko+en",
                confidence=0.99,
            )
        ]


class _CopyInpainter(Inpainter):
    def inpaint(self, slide_image_path: Path, mask_path: Path, output_path: Path) -> Path:
        output_path.write_bytes(slide_image_path.read_bytes())
        return output_path


def test_pipeline_run_creates_artifacts_and_editable_ppt(tmp_path) -> None:
    slide_path = tmp_path / "slide.png"
    Image.new("RGB", (1200, 800), "white").save(slide_path)

    config = PipelineConfig()
    pipeline = SlideToPptPipeline(
        config=config,
        components=PipelineComponents(
            detector=_FakeDetector(),
            ocr_engine=_FakeOcrEngine(),
            inpainter=_CopyInpainter(),
            typesetter=FixedFontTypesetter(font_policy=config.fonts, dpi=config.default_dpi),
            renderer=PowerPointRenderer(),
        ),
    )

    paths = pipeline.run(slide_path, SlideSize(width_px=1200, height_px=800), tmp_path / "artifacts")

    assert paths.detections_json.exists()
    assert paths.ocr_json.exists()
    assert paths.mask_png.exists()
    assert paths.background_png.exists()
    assert paths.result_pptx.exists()

    detections_payload = json.loads(paths.detections_json.read_text(encoding="utf-8"))
    ocr_payload = json.loads(paths.ocr_json.read_text(encoding="utf-8"))
    assert detections_payload["backend"] == "fake_detector"
    assert detections_payload["pipeline"] == {
        "detector_backend": "ocr",
        "ocr_backend": "portable",
        "inpainting_backend": "telea",
    }
    assert detections_payload["runtime"]["platform"]
    assert detections_payload["detections"][0]["box_id"] == "t0001"
    assert ocr_payload["backend"] == "fake_ocr"
    assert ocr_payload["pipeline"] == {
        "detector_backend": "ocr",
        "ocr_backend": "portable",
        "inpainting_backend": "telea",
    }
    assert ocr_payload["runtime"]["python_version"]
    assert ocr_payload["results"][0]["text"] == "Hello PPT"

    presentation = Presentation(str(paths.result_pptx))
    assert len(presentation.slides) == 1
    texts = [shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")]
    assert "Hello PPT" in texts
