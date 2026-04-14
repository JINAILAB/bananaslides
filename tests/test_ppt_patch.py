from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from bananaslides.utils.ppt_patch import inspect_ppt_text, patch_ppt_preserve_style


def test_patch_ppt_preserve_style_keeps_font_and_size(tmp_path: Path) -> None:
    input_pptx = tmp_path / "input.pptx"
    output_pptx = tmp_path / "output.pptx"

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(0, 0, 1000000, 300000)
    paragraph = shape.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = "Original Title"
    run.font.name = "Pretendard"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("123456")
    presentation.save(input_pptx)

    patch_ppt_preserve_style(
        input_pptx,
        output_pptx,
        {
            "slide_number": 1,
            "mappings": [
                {
                    "shape_index": 0,
                    "corrected_text": "Corrected Title",
                    "approx_hex": "#F7F7F7",
                }
            ],
        },
    )

    patched = Presentation(output_pptx)
    patched_run = patched.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]

    assert patched_run.text == "Corrected Title"
    assert patched_run.font.name == "Pretendard"
    assert patched_run.font.size.pt == 24
    assert patched_run.font.bold is True
    assert str(patched_run.font.color.rgb) == "F7F7F7"


def test_inspect_ppt_text_returns_shape_inventory(tmp_path: Path) -> None:
    input_pptx = tmp_path / "inventory.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(100, 200, 300, 400)
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Inventory Text"
    run.font.name = "Pretendard"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor.from_string("0A0B0C")
    presentation.save(input_pptx)

    payload = inspect_ppt_text(input_pptx, slide_number=1)

    item = payload["slides"][0]["items"][0]
    assert item["shape_index"] == 0
    assert item["text"] == "Inventory Text"
    assert item["font_name"] == "Pretendard"
    assert item["font_size_pt"] == 18.0
    assert item["color_hex"] == "#0A0B0C"
