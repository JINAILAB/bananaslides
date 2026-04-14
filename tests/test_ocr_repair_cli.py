from __future__ import annotations

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from typer.testing import CliRunner

from bananaslides import cli as cli_module
from bananaslides.cli import app
from bananaslides.config import ModelAssetPaths
from bananaslides.domain.models import DetectionBox, ElementCategory, OCRLine, OCRResult, RGBColor, SlideSize, TextPlacement
from bananaslides.modules.inpainting.base import Inpainter
from bananaslides.modules.ppt.render import PowerPointRenderer
from bananaslides.pipeline.orchestrator import PipelineComponents
from bananaslides.utils.artifacts import (
    read_ocr_artifact,
    write_detection_artifact,
    write_ocr_artifact,
)
from bananaslides.utils.correction import correct_ocr_results, load_expected_texts_from_deck_plan
from bananaslides.utils.ocr_models import ResolvedOcrModelAssets
from bananaslides.utils.pdf import RenderedPdfPage


def test_load_expected_texts_and_correct_ocr_results(tmp_path: Path) -> None:
    deck_plan = tmp_path / "deck-plan.json"
    deck_plan.write_text(
        json.dumps(
            {
                "slides": [
                    {
                        "slide_number": 1,
                        "title": "Why cuffless blood pressure still matters",
                        "subtitle": "1.4B adults live with hypertension",
                        "on_slide_copy": ["1.4B with hypertension", "23% controlled"],
                        "quantified_points": ["600 million unaware"],
                    }
                ]
            }
        ),
        encoding="utf-8",
    )

    expected = load_expected_texts_from_deck_plan(deck_plan, 1)
    assert expected[0] == "Why cuffless blood pressure still matters"

    corrected, corrections = correct_ocr_results(
        [
            OCRResult(
                box_id="t1",
                text="Why cuffiess blood pressure still matters",
                lines=[OCRLine(text="Why cuffiess blood pressure still matters")],
            )
        ],
        expected,
        min_score=0.7,
    )

    assert corrections[0].applied is True
    assert corrected[0].text == "Why cuffless blood pressure still matters"
    assert corrected[0].lines[0].text == "Why cuffless blood pressure still matters"


def test_repair_ocr_and_render_from_artifacts_cli(tmp_path: Path) -> None:
    runner = CliRunner()
    slide_path = tmp_path / "slide.png"
    background_path = tmp_path / "background.png"
    Image.new("RGB", (1200, 800), "white").save(slide_path)
    Image.new("RGB", (1200, 800), "white").save(background_path)

    detections_json = tmp_path / "slide.detections.json"
    ocr_json = tmp_path / "slide.ocr.json"
    pptx_path = tmp_path / "slide.pptx"
    correction_report = tmp_path / "slide.ocr.corrections.json"
    deck_plan = tmp_path / "deck-plan.json"

    deck_plan.write_text(
        json.dumps(
            {
                "slides": [
                    {
                        "slide_number": 2,
                        "title": "Why cuffless blood pressure still matters",
                        "subtitle": "1.4B adults live with hypertension",
                        "on_slide_copy": ["1.4B with hypertension"],
                        "quantified_points": ["23% controlled"],
                    }
                ]
            }
        ),
        encoding="utf-8",
    )

    write_detection_artifact(
        detections_json,
        image_path=slide_path,
        slide_size=SlideSize(width_px=1200, height_px=800),
        backend="fake_detector",
        detections=[
            DetectionBox(
                box_id="t1",
                category=ElementCategory.TEXT,
                x=100,
                y=100,
                width=500,
                height=80,
            )
        ],
    )
    write_ocr_artifact(
        ocr_json,
        image_path=slide_path,
        backend="fake_ocr",
        language="en",
        results=[
            OCRResult(
                box_id="t1",
                text="Why cuffiess blood pressure still matters",
                lines=[
                    OCRLine(
                        text="Why cuffiess blood pressure still matters",
                        bbox=[[100.0, 100.0], [600.0, 100.0], [600.0, 150.0], [100.0, 150.0]],
                    )
                ],
            )
        ],
    )

    result = runner.invoke(
        app,
        [
            "repair-ocr",
            str(ocr_json),
            "--correction-json",
            str(correction_report),
            "--deck-plan-file",
            str(deck_plan),
            "--slide-number",
            "2",
        ],
    )
    assert result.exit_code == 0, result.stdout
    corrected_results = read_ocr_artifact(ocr_json)
    assert corrected_results[0].text == "Why cuffless blood pressure still matters"
    assert correction_report.exists()
    corrected_payload = json.loads(ocr_json.read_text(encoding="utf-8"))
    assert corrected_payload["backend"] == "ocr-correction"
    assert corrected_payload["runtime"]["platform"]
    assert "pipeline" not in corrected_payload

    result = runner.invoke(
        app,
        [
            "render-from-artifacts",
            str(detections_json),
            str(ocr_json),
            str(background_path),
            str(pptx_path),
        ],
    )
    assert result.exit_code == 0, result.stdout
    assert pptx_path.exists()

    presentation = Presentation(str(pptx_path))
    texts = [shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")]
    assert "Why cuffless blood pressure still matters" in texts


class _FakeDetector:
    backend_name = "fake_detector"

    def detect(self, slide_image: Path, slide_size: SlideSize) -> list[DetectionBox]:
        assert slide_image.exists()
        assert slide_size.width_px == 1200
        return [
            DetectionBox(
                box_id="t1",
                category=ElementCategory.TEXT,
                x=100,
                y=100,
                width=500,
                height=80,
            )
        ]


class _FakeOcrEngine:
    backend_name = "fake_ocr"

    def recognize(self, slide_image: Path, detections: list[DetectionBox]) -> list[OCRResult]:
        assert slide_image.exists()
        assert detections[0].box_id == "t1"
        return [
            OCRResult(
                box_id="t1",
                text="Portable OCR",
                lines=[
                    OCRLine(
                        text="Portable OCR",
                        bbox=[[100.0, 100.0], [400.0, 100.0], [400.0, 150.0], [100.0, 150.0]],
                    )
                ],
            )
        ]


def test_detect_and_ocr_cli_write_runtime_and_pipeline_metadata(tmp_path: Path, monkeypatch) -> None:
    runner = CliRunner()
    slide_path = tmp_path / "slide.png"
    Image.new("RGB", (1200, 800), "white").save(slide_path)

    detections_json = tmp_path / "slide.detections.json"
    ocr_json = tmp_path / "slide.ocr.json"

    monkeypatch.setattr(cli_module, "build_text_detector", lambda config: _FakeDetector())
    monkeypatch.setattr(cli_module, "build_ocr_engine", lambda config: _FakeOcrEngine())

    result = runner.invoke(
        app,
        [
            "detect-text",
            str(slide_path),
            "--output-json",
            str(detections_json),
            "--width-px",
            "1200",
            "--height-px",
            "800",
        ],
    )
    assert result.exit_code == 0, result.stdout

    detections_payload = json.loads(detections_json.read_text(encoding="utf-8"))
    assert detections_payload["backend"] == "fake_detector"
    assert detections_payload["pipeline"] == {
        "detector_backend": "ocr",
        "ocr_backend": "portable",
        "inpainting_backend": "telea",
    }
    assert detections_payload["runtime"]["platform"]

    result = runner.invoke(
        app,
        [
            "ocr-text",
            str(slide_path),
            str(detections_json),
            "--output-json",
            str(ocr_json),
        ],
    )
    assert result.exit_code == 0, result.stdout

    ocr_payload = json.loads(ocr_json.read_text(encoding="utf-8"))
    assert ocr_payload["backend"] == "fake_ocr"
    assert ocr_payload["pipeline"] == {
        "detector_backend": "ocr",
        "ocr_backend": "portable",
        "inpainting_backend": "telea",
    }
    assert ocr_payload["runtime"]["python_version"]


class _DeckFakeDetector:
    backend_name = "deck_fake_detector"

    def detect(self, slide_image_path: Path, slide_size: SlideSize) -> list[DetectionBox]:
        return [
            DetectionBox(
                box_id="t1",
                category=ElementCategory.TEXT,
                x=100,
                y=50,
                width=200,
                height=40,
            )
        ]


class _DeckFakeOcrEngine:
    backend_name = "deck_fake_ocr"

    def recognize(self, slide_image_path: Path, detections: list[DetectionBox]) -> list[OCRResult]:
        return [
            OCRResult(
                box_id="t1",
                text=slide_image_path.stem,
                lines=[
                    OCRLine(
                        text=slide_image_path.stem,
                        bbox=[[100.0, 50.0], [300.0, 50.0], [300.0, 90.0], [100.0, 90.0]],
                    )
                ],
            )
        ]


class _DeckFakeInpainter(Inpainter):
    def inpaint(self, slide_image_path: Path, mask_path: Path, output_path: Path) -> Path:
        output_path.write_bytes(slide_image_path.read_bytes())
        return output_path


class _DeckFakeTypesetter:
    def build_text_placements(
        self,
        detections,
        ocr_results,
        *,
        source_image_path=None,
        background_image_path=None,
    ) -> list[TextPlacement]:
        return [
            TextPlacement(
                box_id="t1",
                text=ocr_results[0].text,
                x=100,
                y=50,
                width=200,
                height=40,
                font_name="Pretendard",
                font_size_pt=20,
                color=RGBColor(0, 0, 0),
            )
        ]


def test_deck_cli_builds_multi_slide_presentation_in_input_order(tmp_path: Path, monkeypatch) -> None:
    runner = CliRunner()
    first_slide = tmp_path / "slide-b.png"
    second_slide = tmp_path / "slide-a.png"
    Image.new("RGB", (1200, 800), "white").save(first_slide)
    Image.new("RGB", (600, 400), "white").save(second_slide)
    output_dir = tmp_path / "deck-artifacts"

    monkeypatch.setattr(
        cli_module,
        "build_default_components",
        lambda config: PipelineComponents(
            detector=_DeckFakeDetector(),
            ocr_engine=_DeckFakeOcrEngine(),
            inpainter=_DeckFakeInpainter(),
            typesetter=_DeckFakeTypesetter(),
            renderer=PowerPointRenderer(),
        ),
    )

    result = runner.invoke(
        app,
        [
            "deck",
            str(first_slide),
            str(second_slide),
            "--output-dir",
            str(output_dir),
        ],
    )
    assert result.exit_code == 0, result.stdout

    deck_pptx = output_dir / "deck.pptx"
    slide_01_pptx = output_dir / "slide-01" / "slide-b.pptx"
    slide_02_pptx = output_dir / "slide-02" / "slide-a.pptx"
    assert deck_pptx.exists()
    assert slide_01_pptx.exists()
    assert slide_02_pptx.exists()
    assert (output_dir / "slide-01" / "slide-b.detections.json").exists()
    assert (output_dir / "slide-02" / "slide-a.detections.json").exists()

    presentation = Presentation(str(deck_pptx))
    assert len(presentation.slides) == 2
    assert presentation.slide_width == 11430000
    assert presentation.slide_height == 7620000

    first_slide_text_shapes = [shape for shape in presentation.slides[0].shapes if hasattr(shape, "text")]
    second_slide_text_shapes = [shape for shape in presentation.slides[1].shapes if hasattr(shape, "text")]
    assert first_slide_text_shapes[0].text == "slide-b"
    assert second_slide_text_shapes[0].text == "slide-a"
    assert first_slide_text_shapes[0].left == 952500
    assert second_slide_text_shapes[0].left == 1905000


def test_run_cli_rejects_pdf_input(tmp_path: Path) -> None:
    runner = CliRunner()
    pdf_path = tmp_path / "slides.pdf"
    pdf_path.write_bytes(b"%PDF-1.7\n")

    result = runner.invoke(app, ["run", str(pdf_path)])

    assert result.exit_code != 0
    output = result.stdout + result.stderr
    assert "only supports raster slide images" in output
    assert "bananaslides deck" in output


def test_deck_cli_expands_pdf_pages(tmp_path: Path, monkeypatch) -> None:
    runner = CliRunner()
    pdf_path = tmp_path / "brief.pdf"
    pdf_path.write_bytes(b"%PDF-1.7\n")
    page_1 = tmp_path / "brief-page-001.png"
    page_2 = tmp_path / "brief-page-002.png"
    Image.new("RGB", (1200, 800), "white").save(page_1)
    Image.new("RGB", (1200, 800), "white").save(page_2)
    output_dir = tmp_path / "deck-from-pdf"

    monkeypatch.setattr(
        cli_module,
        "build_default_components",
        lambda config: PipelineComponents(
            detector=_DeckFakeDetector(),
            ocr_engine=_DeckFakeOcrEngine(),
            inpainter=_DeckFakeInpainter(),
            typesetter=_DeckFakeTypesetter(),
            renderer=PowerPointRenderer(),
        ),
    )
    monkeypatch.setattr(
        cli_module,
        "render_pdf_pages",
        lambda pdf_path, output_dir, dpi=144: [
            RenderedPdfPage(
                source_pdf_path=pdf_path,
                page_number=1,
                image_path=page_1,
                slide_size=SlideSize(width_px=1200, height_px=800),
            ),
            RenderedPdfPage(
                source_pdf_path=pdf_path,
                page_number=2,
                image_path=page_2,
                slide_size=SlideSize(width_px=1200, height_px=800),
            ),
        ],
    )

    result = runner.invoke(
        app,
        [
            "deck",
            str(pdf_path),
            "--output-dir",
            str(output_dir),
        ],
    )
    assert result.exit_code == 0, result.stdout

    deck_pptx = output_dir / "deck.pptx"
    assert deck_pptx.exists()
    assert (output_dir / "slide-01" / "brief-page-001.detections.json").exists()
    assert (output_dir / "slide-02" / "brief-page-002.detections.json").exists()

    presentation = Presentation(str(deck_pptx))
    assert len(presentation.slides) == 2
    first_slide_text_shapes = [shape for shape in presentation.slides[0].shapes if hasattr(shape, "text")]
    second_slide_text_shapes = [shape for shape in presentation.slides[1].shapes if hasattr(shape, "text")]
    assert first_slide_text_shapes[0].text == "brief-page-001"
    assert second_slide_text_shapes[0].text == "brief-page-002"


def test_list_ocr_presets_cli_includes_default_preset() -> None:
    runner = CliRunner()

    result = runner.invoke(app, ["list-ocr-presets"])

    assert result.exit_code == 0, result.stdout
    assert "ko-en" in result.stdout
    assert "Korean + English" in result.stdout


def test_init_models_cli_uses_shared_installer(tmp_path: Path, monkeypatch) -> None:
    runner = CliRunner()
    installed: dict[str, object] = {}

    def fake_install(preset: str, model_home: Path, *, activate: bool, force: bool) -> ResolvedOcrModelAssets:
        model_home.mkdir(parents=True, exist_ok=True)
        assets = ModelAssetPaths(
            rapidocr_det_model=model_home / "presets" / preset / "det.onnx",
            rapidocr_cls_model=model_home / "presets" / preset / "cls.onnx",
            rapidocr_rec_model=model_home / "presets" / preset / "rec.onnx",
            rapidocr_keys=model_home / "presets" / preset / "dict.txt",
        )
        assets.rapidocr_det_model.parent.mkdir(parents=True, exist_ok=True)
        for path in (
            assets.rapidocr_det_model,
            assets.rapidocr_cls_model,
            assets.rapidocr_rec_model,
            assets.rapidocr_keys,
        ):
            path.write_bytes(b"stub")
        installed.update({"preset": preset, "model_home": model_home, "activate": activate, "force": force})
        return ResolvedOcrModelAssets(
            model_assets=assets,
            source="preset",
            preset_id=preset,
            model_home=model_home,
        )

    monkeypatch.setattr(cli_module, "install_ocr_preset", fake_install)

    result = runner.invoke(
        app,
        [
            "init-models",
            "--preset",
            "latin",
            "--ocr-model-home",
            str(tmp_path / "models"),
        ],
    )

    assert result.exit_code == 0, result.stdout
    assert installed["preset"] == "latin"
    assert installed["activate"] is True
    assert "Installed OCR preset: latin" in result.stdout
