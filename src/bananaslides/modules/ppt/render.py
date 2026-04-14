from __future__ import annotations

from pathlib import Path

from bananaslides.domain.models import PresentationRenderSpec
from bananaslides.utils.geometry import pixels_to_emu


class PowerPointRenderer:
    """Render background images and text placements into an editable PPT."""

    def render(self, spec: PresentationRenderSpec, output_path: Path) -> Path:
        from pptx import Presentation
        from pptx.dml.color import RGBColor as PptxRGBColor
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Pt

        output_path.parent.mkdir(parents=True, exist_ok=True)

        presentation = Presentation()
        first_slide = spec.slides[0]
        presentation.slide_width = pixels_to_emu(first_slide.slide_size.width_px, dpi=spec.dpi)
        presentation.slide_height = pixels_to_emu(first_slide.slide_size.height_px, dpi=spec.dpi)
        blank_layout = presentation.slide_layouts[6]

        for slide_spec in spec.slides:
            slide = presentation.slides.add_slide(blank_layout)

            if slide_spec.background_image_path is not None:
                slide.shapes.add_picture(
                    str(slide_spec.background_image_path),
                    0,
                    0,
                    width=presentation.slide_width,
                    height=presentation.slide_height,
                )

            for placement in slide_spec.image_placements:
                slide.shapes.add_picture(
                    str(placement.image_path),
                    pixels_to_emu(placement.x, spec.dpi),
                    pixels_to_emu(placement.y, spec.dpi),
                    width=pixels_to_emu(placement.width, spec.dpi),
                    height=pixels_to_emu(placement.height, spec.dpi),
                )

            for placement in slide_spec.text_placements:
                shape = slide.shapes.add_textbox(
                    pixels_to_emu(placement.x, spec.dpi),
                    pixels_to_emu(placement.y, spec.dpi),
                    pixels_to_emu(placement.width, spec.dpi),
                    pixels_to_emu(placement.height, spec.dpi),
                )
                text_frame = shape.text_frame
                text_frame.word_wrap = placement.word_wrap
                text_frame.margin_left = 0
                text_frame.margin_right = 0
                text_frame.margin_top = 0
                text_frame.margin_bottom = 0
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
                text_frame.text = placement.text

                for paragraph in text_frame.paragraphs:
                    paragraph.alignment = self._map_alignment(placement.align, PP_ALIGN)
                    runs = paragraph.runs or [paragraph.add_run()]
                    for run in runs:
                        run.font.name = placement.font_name
                        run.font.size = Pt(placement.font_size_pt)
                        run.font.bold = placement.bold
                        run.font.italic = placement.italic
                        run.font.color.rgb = PptxRGBColor(*placement.color.as_tuple())

                if placement.auto_fit:
                    text_frame.fit_text(
                        font_family=placement.font_name,
                        max_size=max(1, int(round(placement.font_size_pt))),
                        bold=placement.bold,
                        italic=placement.italic,
                        font_file=placement.font_file,
                    )

        presentation.save(output_path)
        return output_path

    @staticmethod
    def _map_alignment(align: str, align_enum: object) -> object:
        mapping = {
            "left": align_enum.LEFT,
            "center": align_enum.CENTER,
            "right": align_enum.RIGHT,
        }
        return mapping.get(align.lower(), align_enum.LEFT)
