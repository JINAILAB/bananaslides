from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor


def inspect_ppt_text(pptx_path: Path, slide_number: int | None = None) -> dict[str, Any]:
    presentation = Presentation(pptx_path)
    slide_indices = [slide_number - 1] if slide_number is not None else list(range(len(presentation.slides)))

    slides: list[dict[str, Any]] = []
    for slide_index in slide_indices:
        slide = presentation.slides[slide_index]
        items: list[dict[str, Any]] = []
        for shape_index, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if not text:
                continue

            style = _capture_text_frame_style(shape.text_frame)
            items.append(
                {
                    "shape_index": shape_index,
                    "text": text,
                    "left": int(shape.left),
                    "top": int(shape.top),
                    "width": int(shape.width),
                    "height": int(shape.height),
                    "font_name": style["font_name"],
                    "font_size_pt": style["font_size_pt"],
                    "bold": style["bold"],
                    "italic": style["italic"],
                    "alignment": style["alignment"],
                    "color_hex": style["color_hex"],
                }
            )
        slides.append({"slide_number": slide_index + 1, "items": items})

    return {
        "pptx_path": str(pptx_path),
        "slides": slides,
    }


def save_ppt_text_inventory(pptx_path: Path, output_json: Path, slide_number: int | None = None) -> Path:
    payload = inspect_ppt_text(pptx_path, slide_number=slide_number)
    output_json.parent.mkdir(parents=True, exist_ok=True)
    output_json.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    return output_json


def patch_ppt_preserve_style(
    input_pptx: Path,
    output_pptx: Path,
    mapping_payload: dict[str, Any],
    *,
    slide_number: int | None = None,
) -> Path:
    presentation = Presentation(input_pptx)
    patches = _resolve_patches(mapping_payload, slide_number=slide_number)

    for slide_idx, mappings in patches.items():
        slide = presentation.slides[slide_idx]
        used_shape_indexes: set[int] = set()
        for mapping in mappings:
            shape_index = _resolve_shape_index(slide, mapping, used_shape_indexes)
            if shape_index is None:
                continue
            used_shape_indexes.add(shape_index)
            shape = slide.shapes[shape_index]
            if not getattr(shape, "has_text_frame", False):
                continue
            _apply_mapping_to_shape(shape.text_frame, mapping)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_pptx)
    return output_pptx


def _resolve_patches(mapping_payload: dict[str, Any], *, slide_number: int | None) -> dict[int, list[dict[str, Any]]]:
    if "slides" in mapping_payload:
        return {
            int(item["slide_number"]) - 1: list(item.get("mappings", []))
            for item in mapping_payload.get("slides", [])
        }

    mappings = list(mapping_payload.get("mappings", []))
    resolved_slide_number = slide_number or mapping_payload.get("slide_number") or 1
    return {int(resolved_slide_number) - 1: mappings}


def _resolve_shape_index(slide: Any, mapping: dict[str, Any], used_shape_indexes: set[int]) -> int | None:
    if mapping.get("shape_index") is not None:
        return int(mapping["shape_index"])

    current_text = str(mapping.get("current_text", "")).strip()
    if not current_text:
        return None

    for shape_index, shape in enumerate(slide.shapes):
        if shape_index in used_shape_indexes:
            continue
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape.text.strip() == current_text:
            return shape_index
    return None


def _apply_mapping_to_shape(text_frame: Any, mapping: dict[str, Any]) -> None:
    style = _capture_text_frame_style(text_frame)
    new_text = mapping.get("corrected_text")
    if new_text is None:
        new_text = text_frame.text

    _rewrite_text_frame(text_frame, str(new_text), style)

    color_hex = mapping.get("approx_hex")
    if not color_hex:
        return
    rgb = _hex_to_rgb(str(color_hex))
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = rgb


def _capture_text_frame_style(text_frame: Any) -> dict[str, Any]:
    paragraph0 = text_frame.paragraphs[0]
    run0 = None
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text is not None:
                run0 = run
                break
        if run0 is not None:
            break

    color_hex = None
    if run0 is not None and run0.font.color is not None and run0.font.color.rgb is not None:
        color_hex = f"#{run0.font.color.rgb}"

    return {
        "font_name": run0.font.name if run0 is not None else None,
        "font_size": run0.font.size if run0 is not None else None,
        "font_size_pt": float(run0.font.size.pt) if run0 is not None and run0.font.size else None,
        "bold": run0.font.bold if run0 is not None else None,
        "italic": run0.font.italic if run0 is not None else None,
        "alignment": paragraph0.alignment,
        "word_wrap": text_frame.word_wrap,
        "vertical_anchor": text_frame.vertical_anchor,
        "margin_left": text_frame.margin_left,
        "margin_right": text_frame.margin_right,
        "margin_top": text_frame.margin_top,
        "margin_bottom": text_frame.margin_bottom,
        "color_hex": color_hex,
    }


def _rewrite_text_frame(text_frame: Any, text: str, style: dict[str, Any]) -> None:
    text_frame.clear()
    text_frame.word_wrap = style["word_wrap"]
    text_frame.vertical_anchor = style["vertical_anchor"]
    text_frame.margin_left = style["margin_left"]
    text_frame.margin_right = style["margin_right"]
    text_frame.margin_top = style["margin_top"]
    text_frame.margin_bottom = style["margin_bottom"]

    lines = text.splitlines() or [text]
    for index, line in enumerate(lines):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = style["alignment"]
        run = paragraph.add_run()
        run.text = line
        if style["font_name"] is not None:
            run.font.name = style["font_name"]
        if style["font_size"] is not None:
            run.font.size = style["font_size"]
        run.font.bold = style["bold"]
        run.font.italic = style["italic"]
        if style["color_hex"]:
            run.font.color.rgb = _hex_to_rgb(style["color_hex"])


def _hex_to_rgb(value: str) -> PptxRGBColor:
    normalized = value.strip().lstrip("#")
    return PptxRGBColor.from_string(normalized.upper())
